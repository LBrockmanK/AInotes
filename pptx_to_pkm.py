import os
import shutil
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from transformers import AutoTokenizer, AutoModel
from transformers import pipeline
import torch
import numpy as np
from pptx.shapes.picture import Picture
import errno
import re
from io import BytesIO
import win32com.client
import string
import unicodedata

# Initialize the tokenizer and model for sentence embeddings
tokenizer = AutoTokenizer.from_pretrained("t5-small", model_max_length=512)
model = AutoModel.from_pretrained("sentence-transformers/paraphrase-MiniLM-L6-v2")
imgslide = 0

# Initialize the summarization pipeline
summarizer = pipeline("summarization", model="t5-small")

# Function to compute text similarity using sentence embeddings
def similarity(text1, text2):
    # Keep only alphabetical characters and spaces
    text1 = ''.join(c for c in text1 if c.isalpha() or c == ' ')
    text2 = ''.join(c for c in text2 if c.isalpha() or c == ' ')

    inputs = tokenizer([text1, text2], return_tensors="pt", padding=True, truncation=True)

    # Filter out unknown tokens
    filtered_input_ids = [
        inputs["input_ids"][i, inputs["input_ids"][i] < model.config.vocab_size] for i in range(inputs["input_ids"].shape[0])
    ]

    # Find the maximum length after filtering
    max_length = max([len(ids) for ids in filtered_input_ids])

    # Pad the filtered input tensors to the same length
    padded_input_ids = torch.stack([torch.cat((ids, torch.zeros(max_length - len(ids), dtype=torch.long))) for ids in filtered_input_ids])

    inputs["input_ids"] = padded_input_ids

    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1).numpy()

    similarity_score = np.inner(embeddings[0], embeddings[1]) / (
        np.linalg.norm(embeddings[0]) * np.linalg.norm(embeddings[1]))
    return similarity_score

# Find likely links between subject
def link_slides(content, similarity_threshold):
    # We're going to want to take each title line, and search the text of other slides for similarities and link those back to the working slide
    # Link every slide to its predecessor and follower
    linked_content = []

    for index, slide in enumerate(content):
        print(f"Linking slide: {index}")
        # Skip processing if the subject starts with "Summary"
        if not slide["title"].startswith("Summary of "):
            # Link slides
            if index > 0:  # There is a previous slide
                prev_slide = content[index - 1]
                slide["links"] += f'\n\nPrev: [[{prev_slide["title"]}|{prev_slide["subject"]}]]'

            if index < len(content) - 1:  # There is a next slide
                next_slide = content[index + 1]
                slide["links"] += f'\nNext: [[{next_slide["title"]}|{next_slide["subject"]}]]'

            slide["links"] += f'\nRelated Content:'

            # Perform similarity scan and link related slides
            for other_index, other_slide in enumerate(content):
                if index != other_index:
                    similarity_score = similarity(slide["subject"], other_slide["text"])
                    if similarity_score > similarity_threshold:
                        other_slide["links"] += f'\n[[{slide["title"]}|{slide["subject"]}]]'

        linked_content.append(slide)

    return linked_content

# Function to process and merge similar slides in a content list
def process_slides(content, length):
    # Summarization and other operations
    processed_content = []
    for index, slide in enumerate(content):
        print(f"Processing slide: {index}")
        # Skip processing if the subject starts with "Summary"
        if not slide["subject"].startswith("Summary"):

            # Summarize the text in the slide
            slide["text"] = summarizer(slide["text"], max_length=length, do_sample=False)[0]["summary_text"]

        processed_content.append(slide)

    return processed_content


def merge_slides(content, similarity_threshold):
    merged_content = []
    # TODO: Do we want to merge the subjects together then set the title based on that?
    # TODO: Maybe instead of unique names, have sub folders for different sources, Will make things read better
    # Merging similar slides
    while content:
        current_slide = content.pop(0)
        merged_slide = current_slide.copy()

        while content:
            next_slide = content[0]

            if similarity(current_slide["subject"], next_slide["subject"]) > similarity_threshold:
                # Merge slides
                print(f"Merging slides: '{current_slide['subject']}' and '{next_slide['subject']}'")
                merged_slide["text"] += "\n" + next_slide["text"]
                merged_slide["images"].extend(next_slide["images"])
                merged_slide["notes"].append(next_slide["notes"][0])

                # Remove the merged slide from the content list
                content.pop(0)
            else:
                print(f"Not merging slides: '{current_slide['subject']}' and '{next_slide['subject']}'")
                break

        # Add the merged_slide to the merged_content list
        merged_content.append(merged_slide)

        # Check if the final slide's subject is similar to "Summary"
    summary_start = None
    for i in range(len(merged_content) - 1, -1, -1):
        if similarity(merged_content[i]["subject"], "Summary") > 0.5:
            summary_start = i
        else:
            break

    if summary_start is not None:
        # Merge all summary slides found
        summary_slide = merged_content[summary_start]
        for slide in merged_content[summary_start + 1:]:
            summary_slide["text"] += "\n" + slide["text"]
            summary_slide["images"].extend(slide["images"])

        # Add a new line and "Content:" to the merged summary slide's text field
        summary_slide["links"] += "\n\nContent:"

        # Add a new line with the subject line enclosed in double square brackets for each object in
        for slide in merged_content[:summary_start]:  # Exclude the summary slides
            summary_slide["links"] += f"\n[[{slide['title']}|{slide['subject']}]]"

        # Set the subject line for the combined summary slide
        summary_slide["title"] = "Summary of " + summary_slide["source"]

        # Remove the merged summary slides from the merged_content list
        merged_content = merged_content[:summary_start]

        # Append the combined summary slide to the end of the merged_content list
        merged_content.append(summary_slide)

    for slide in merged_content:
        slide["text"] = re.sub(r'\n\s*\n', '\n\n', slide["text"])

    return merged_content

# Reset output directory
def clear_directory(directory):
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                shutil.rmtree(file_path)
        except Exception as e:
            print(f"Failed to delete {file_path}. Reason: {e}")

# Function to extract text and images from a PPTX file
def extract_content(pptx_file, output_dir):
    # Read the PPTX file using the python-pptx library and store it in the variable 'prs'.
    prs = Presentation(pptx_file)
    # Initialize an empty list to store the extracted content.
    content = []
    previous_subject = "UNKNOWN"
    source_name = os.path.splitext(os.path.basename(pptx_file))[0]

    # Access PowerPoint Application
    Application = win32com.client.Dispatch("PowerPoint.Application")
    win32_presentation = Application.Presentations.Open(os.path.abspath(pptx_file), WithWindow=False)  # win32com Presentation

    # Iterate through the slides in the presentation using their index and value.
    for slide_num, slide in enumerate(prs.slides):
        # Initialize a dictionary to store the text and images for each slide.
        slide_content = {"subject": "", "text": "", "images": [], "source": source_name, "title": "", "links" : "", "notes" : []}
        text_boxes = []

        # Iterate through the shapes in each slide.
        for shape in slide.shapes:
            # Check if the shape has a text frame and append its text to the text_boxes list.
            if shape.has_text_frame:
                text_boxes.append(shape.text)

        # Add notes text to text_boxes
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if shape.has_text_frame:
                    notes_text += shape.text + "\n"
            text_boxes.append(notes_text.strip())
            
        slide_content["notes"].append(notes_text.strip())

        # Assign the first text box as the subject, and join the rest of the text boxes with a newline.
        if text_boxes:
            # Remove characters that make for bad file names
            slide_content["subject"] = " ".join("".join(c if c.isalpha() or c.isspace() else " " for c in text_boxes[0]).strip().split())

            # Combine separate text fields
            slide_content["text"] = "\n".join(text_boxes[1:])
        else:
            slide_content["subject"] = previous_subject

        # Update the previous_subject variable with the current slide's subject.
        previous_subject = slide_content["subject"]

        slide_content["title"] = slide_content["subject"] + " " + slide_content["source"]

        # Add the whole slide image to the image field
        slide_image_path = os.path.abspath(os.path.join(output_dir, f"{source_name}_slide_{slide_num}.jpg"))
        win32_presentation.Slides[slide_num].Export(slide_image_path, "JPG")
        slide_content["images"].append(os.path.basename(slide_image_path))

        content.append(slide_content)

    # Close PowerPoint Application
    win32_presentation.Close()
    Application.Quit()
    win32_presentation = None
    Application = None

    # Return the content list containing text and images from the PPTX file.
    return content

def is_special_char(c):
    category = unicodedata.category(c)
    return category.startswith("M") or category.startswith("S")

def check_directory_writable(directory):
    try:
        test_file = os.path.join(directory, "test_write.txt")
        with open(test_file, "w") as file:
            file.write("test")
        os.remove(test_file)
        return True
    except IOError as e:
        if e.errno == errno.EACCES:
            return False
        else:
            raise

def create_markdown_files(content, output_dir):
    for slide_num, slide in enumerate(content):
        # Replace invalid characters in the subject line for creating a file
        md_filename = f"{slide['title']}.md"
        md_filepath = os.path.join(output_dir, md_filename)

        with open(md_filepath, "w", encoding="utf-8-sig") as md_file:
            if slide["images"] and slide["notes"]:
                md_file.write("- [ ] TODO:\n")
                for image, note in zip(slide["images"], slide["notes"]):
                    md_file.write(f"![[{image}]]\n")
                    md_file.write(f"{note}\n")
                    md_file.write("\n")

            md_file.write(slide["links"])

def main():
    # Set the input directory path where the PPTX files are stored.
    input_directory = "Inputs"
    base_output_dir = "pptx_to_obsidian"
    os.makedirs(base_output_dir, exist_ok=True)
    figures_output_dir = os.path.join(base_output_dir, "Figures")
    os.makedirs(figures_output_dir, exist_ok=True)
    clear_directory(figures_output_dir)
    notes_output_dir = os.path.join(base_output_dir, "Notes")
    os.makedirs(notes_output_dir, exist_ok=True)
    clear_directory(notes_output_dir)

    # Check if the figures_output_dir is writable
    if not check_directory_writable(figures_output_dir):
        print(f"Error: The directory '{figures_output_dir}' is not writable. Aborting.")
        return

    # Check if the notes_output_dir is writable
    if not check_directory_writable(notes_output_dir):
        print(f"Error: The directory '{notes_output_dir}' is not writable. Aborting.")
        return

    # Initialize an empty list to store the content extracted from each PPTX file.
    contents = []

    # Iterate through the files in the input directory using os.walk.
    for root, dirs, files in os.walk(input_directory):
        for file in files:
            # Check if the file has a '.pptx' extension (case-insensitive).
            if file.lower().endswith(".pptx"):
                # Create the full path of the PPTX file by joining the root directory and the file name.
                pptx_path = os.path.join(root, file)
                # Extract the content (text and images) from the PPTX file using the extract_content function.
                content = extract_content(pptx_path, figures_output_dir)
                # Merge slides based on subject similarity
                content = merge_slides(content, 0.6)
                # Append the extracted content to the 'contents' list.
                contents.append(content)

    # Combine all slide decks into one
    contents = [slide for slide_deck in contents for slide in slide_deck]

    # Summarize slide content
    contents = process_slides(contents, 100)

    # Find likely links between subjects and other note contents
    contents = link_slides(contents,0.5)

    # Save slide contents into markdown files
    create_markdown_files(contents, notes_output_dir)

if __name__ == "__main__":
    main()
    print("Done")